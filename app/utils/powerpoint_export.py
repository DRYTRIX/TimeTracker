"""
PowerPoint export utilities for reports
Requires: python-pptx
Install: pip install python-pptx
"""

import io
from datetime import datetime

# Try to import pptx, but make it optional
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # Define dummy classes to prevent errors if module is accessed
    Presentation = None
    Inches = None
    Pt = None
    PP_ALIGN = None
    RGBColor = None


def create_report_powerpoint(entries, title="TimeTracker Report", filename_prefix="timetracker_report"):
    """Create PowerPoint presentation from time entries

    Args:
        entries: List of TimeEntry objects
        title: Presentation title
        filename_prefix: Prefix for the filename

    Returns:
        tuple: (BytesIO object with PPTX file, filename)

    Raises:
        ImportError: If python-pptx is not installed
    """
    if not PPTX_AVAILABLE:
        raise ImportError("PowerPoint export requires python-pptx. Install it with: pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]

    title_shape.text = title
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{len(entries)} time entries"

    # Summary slide
    summary_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(summary_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Summary"

    # Calculate summary
    total_hours = sum(entry.duration_hours for entry in entries if entry.end_time)
    billable_hours = sum(entry.duration_hours for entry in entries if entry.billable and entry.end_time)

    projects_count = len(set(entry.project_id for entry in entries))
    users_count = len(set(entry.user_id for entry in entries))

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Total Hours: {total_hours:.2f}"

    p = tf.add_paragraph()
    p.text = f"Billable Hours: {billable_hours:.2f}"
    p.level = 0

    p = tf.add_paragraph()
    p.text = f"Projects: {projects_count}"
    p.level = 0

    p = tf.add_paragraph()
    p.text = f"Users: {users_count}"
    p.level = 0

    # Time entries slide with table
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Time Entries"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    # Create table (limit to 20 entries per slide for readability)
    rows = min(len(entries), 20) + 1  # +1 for header
    cols = 6

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(0.5)  # ID
    table.columns[1].width = Inches(1.5)  # User
    table.columns[2].width = Inches(1.5)  # Project
    table.columns[3].width = Inches(1.2)  # Date
    table.columns[4].width = Inches(1.0)  # Duration
    table.columns[5].width = Inches(3.3)  # Notes

    # Header row
    headers = ["ID", "User", "Project", "Date", "Hours", "Notes"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Blue
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, entry in enumerate(entries[:20], 1):
        data = [
            str(entry.id),
            entry.user.display_name if entry.user else "Unknown",
            entry.project.name if entry.project else "N/A",
            entry.start_time.strftime("%Y-%m-%d") if entry.start_time else "",
            f"{entry.duration_hours:.2f}" if entry.end_time else "In Progress",
            (entry.notes or "")[:50],  # Truncate long notes
        ]

        for col_idx, value in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray

    # Create additional slides if more than 20 entries
    if len(entries) > 20:
        for i in range(20, len(entries), 20):
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add title
            txBox = slide.shapes.add_textbox(left, top - Inches(1), width, height)
            tf = txBox.text_frame
            tf.text = f"Time Entries (continued) - Page {i // 20 + 2}"
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True

            # Create table for this batch
            batch_entries = entries[i : i + 20]
            rows = len(batch_entries) + 1

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Set column widths
            for col_idx in range(cols):
                table.columns[col_idx].width = table.columns[col_idx].width

            # Header row (same as before)
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Data rows
            for row_idx, entry in enumerate(batch_entries, 1):
                data = [
                    str(entry.id),
                    entry.user.display_name if entry.user else "Unknown",
                    entry.project.name if entry.project else "N/A",
                    entry.start_time.strftime("%Y-%m-%d") if entry.start_time else "",
                    f"{entry.duration_hours:.2f}" if entry.end_time else "In Progress",
                    (entry.notes or "")[:50],
                ]

                for col_idx, value in enumerate(data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value)
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    filename = f"{filename_prefix}_{datetime.now().strftime('%Y%m%d')}.pptx"

    return output, filename
